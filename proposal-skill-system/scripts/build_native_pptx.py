"""
extracted/slide_*.json からネイティブPPTXを生成する。

戦略:
- 1280px × 720px → 13.33in × 7.5in  (16:9)
- 1px = 13.33/1280 inch ≈ 9525 EMU
- HTML要素: 背景色・罫線・テキストを持つ <div> をPPTXシェイプ+text_frameに
- SVG要素:
  - <rect> → PPTX 矩形シェイプ
  - <text> → 内部に含まれるrectがあれば、そのrectのtext_frameに埋め込み (テキストをオブジェクトに直接書く)
  - <line> → PPTX 直線
  - <polygon> → PPTX 自由形状 (主に矢印)
- テキストとシェイプは一体化（テキストボックスを別レイヤーで重ねない）
"""
import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

EXTRACTED_DIR = Path("/home/claude/extracted")
OUTPUT_PPTX = "/mnt/user-data/outputs/booking_v9_chholdings_19pages.pptx"

# 基準サイズ
SRC_W = 1280  # px
SRC_H = 720
DST_W = Inches(13.33)
DST_H = Inches(7.5)

# 1px → EMU
SCALE_X = DST_W / SRC_W
SCALE_Y = DST_H / SRC_H

def px_to_emu(px, axis="x"):
    """px座標をEMUに変換"""
    return int(px * (SCALE_X if axis == "x" else SCALE_Y))

def parse_color(c):
    """CSS color/hexを (R, G, B) tupleに変換。失敗時はNone"""
    if not c or c in ("none", "transparent"):
        return None
    c = c.strip()
    if c.startswith("#"):
        c = c[1:]
        if len(c) == 3:
            c = "".join([x*2 for x in c])
        if len(c) == 6:
            return (int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
    m = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', c)
    if m:
        return (int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return None

def parse_color_with_alpha(c):
    """rgba対応"""
    if not c or c in ("none", "transparent"):
        return None, 1.0
    m = re.match(r'rgba\((\d+),\s*(\d+),\s*(\d+),\s*([\d.]+)', c)
    if m:
        return (int(m.group(1)), int(m.group(2)), int(m.group(3))), float(m.group(4))
    rgb = parse_color(c)
    return rgb, 1.0

def parse_px(s):
    """'12px' → 12.0"""
    if not s:
        return 0.0
    m = re.match(r'(-?[\d.]+)', str(s))
    return float(m.group(1)) if m else 0.0

def font_size_pt(px):
    """px → pt (1pt = 1.333px)"""
    return max(6, round(px / 1.333, 1))

def add_filled_shape(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_w=0, rx=0,
                      shape=MSO_SHAPE.RECTANGLE, dash=False):
    """背景色つき矩形シェイプを追加。返値: shape"""
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    # 塗り
    if fill_rgb is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(*fill_rgb)
    # 枠
    if line_rgb is None or line_w == 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = RGBColor(*line_rgb)
        shp.line.width = Pt(line_w)
        if dash:
            shp.line.dash_style = 7  # MSO_LINE_DASH_STYLE.DASH
    # 角丸
    if rx > 0 and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        # 角丸調整 (adj1)
        try:
            shp.adjustments[0] = min(0.5, rx / min(w/Emu(1), h/Emu(1)) * 914400 / 100000)
        except Exception:
            pass
    # text_frameを取得（シェイプには初期から空のtext_frameがある）
    tf = shp.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    # 既存のparagraphをクリア
    for p in tf.paragraphs:
        p.text = ""
    return shp

def add_textbox(slide, x, y, w, h, text, font_size_px=11, color_rgb=None, weight="400",
                 align="left", v_anchor="top"):
    """シェイプなしのテキストボックスを追加"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    if v_anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif v_anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt(font_size_px))
    run.font.name = "Hiragino Sans"
    if color_rgb:
        run.font.color.rgb = RGBColor(*color_rgb)
    if weight in ("700", "bold") or (weight.isdigit() and int(weight) >= 600):
        run.font.bold = True
    return tb

def write_text_in_shape(shape, texts, anchor="middle"):
    """既存shapeにテキストを書き込む。textsは [(text, font_size_px, color_rgb, weight, align), ...]
    text内の'\\n'は改行paragraphとして扱う"""
    tf = shape.text_frame
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    
    # 既存paragraphsをクリア
    p_xml = tf._txBody
    for p_el in list(p_xml.findall(qn('a:p'))):
        p_xml.remove(p_el)
    
    first = True
    for text, fs, c, w, al in texts:
        # 改行を含む場合は複数paragraphに
        lines = text.split("\n") if text else [""]
        for line in lines:
            p = tf.add_paragraph()
            p.alignment = {"start": PP_ALIGN.LEFT, "left": PP_ALIGN.LEFT,
                            "middle": PP_ALIGN.CENTER, "center": PP_ALIGN.CENTER,
                            "end": PP_ALIGN.RIGHT, "right": PP_ALIGN.RIGHT}.get(al, PP_ALIGN.CENTER)
            run = p.add_run()
            run.text = line.strip()
            run.font.size = Pt(font_size_pt(fs))
            run.font.name = "Hiragino Sans"
            if c:
                run.font.color.rgb = RGBColor(*c)
            if w in ("700", "bold") or (str(w).isdigit() and int(w) >= 600):
                run.font.bold = True

def add_line(slide, x1, y1, x2, y2, color_rgb=(136,136,136), width_pt=0.75, dash=False):
    """直線を追加"""
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = RGBColor(*color_rgb)
    line.line.width = Pt(width_pt)
    if dash:
        line.line.dash_style = 7
    return line

def text_in_rect(text_x, text_y, rect):
    """テキストの座標がrectの内側にあるか判定"""
    return (rect["x"] <= text_x <= rect["x"] + rect["w"] and
            rect["y"] - 5 <= text_y <= rect["y"] + rect["h"] + 5)

def build_slide(prs, slide_data):
    """1枚分のスライドをビルド"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    
    html_elems = slide_data["htmlElements"]
    svg_elems = slide_data["svgElements"]
    
    # === 0. スライドルート背景を塗る ===
    slide_bg_rgb, _ = parse_color_with_alpha(slide_data.get("slideBackground", ""))
    if slide_bg_rgb and slide_bg_rgb != (255, 255, 255):
        bg_shp = add_filled_shape(slide, 0, 0, prs.slide_width, prs.slide_height, fill_rgb=slide_bg_rgb)
    
    # === 1. HTML要素を変換 ===
    # 親→子の順なので、面積の大きい順にソート（背景から先に描画）
    html_elems_sorted = sorted(html_elems, key=lambda e: -(e.get("w", 0) * e.get("h", 0)))
    
    for elem in html_elems_sorted:
        x = px_to_emu(elem["x"], "x")
        y = px_to_emu(elem["y"], "y")
        w = px_to_emu(elem["w"], "x")
        h = px_to_emu(elem["h"], "y")
        if w <= 0 or h <= 0:
            continue
        
        bg_rgb, bg_alpha = parse_color_with_alpha(elem.get("bg", ""))
        bt_w = parse_px(elem.get("borderTopWidth", "0"))
        bl_w = parse_px(elem.get("borderLeftWidth", "0"))
        br = parse_px(elem.get("borderRadius", "0"))
        
        has_bg = bg_rgb is not None and bg_alpha > 0.05
        has_top_border = bt_w > 0
        has_left_border = bl_w > 0
        
        # スライド全体のような大きすぎる要素はスキップ
        if elem.get("cls") == "slide" or (elem["w"] >= 1270 and elem["h"] >= 700):
            continue
        
        # 線色
        line_rgb = None
        line_w_pt = 0
        if has_top_border:
            line_rgb = parse_color(elem.get("borderTopColor", "")) or (200, 200, 200)
            line_w_pt = bt_w
        elif has_left_border:
            lc = parse_color(elem.get("borderLeftColor", "")) or (200, 200, 200)
            add_filled_shape(slide, x, y, px_to_emu(bl_w, "x"), h, fill_rgb=lc)
        
        text = elem.get("text", "")
        color_rgb = parse_color(elem.get("color", "")) or (0, 0, 0)
        fs = parse_px(elem.get("fontSize", "11px"))
        weight = elem.get("fontWeight", "400")
        align = elem.get("textAlign", "left")
        
        # シェイプ作成
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if br > 2 else MSO_SHAPE.RECTANGLE
        if has_bg or line_rgb:
            shp = add_filled_shape(slide, x, y, w, h,
                                     fill_rgb=bg_rgb if has_bg else None,
                                     line_rgb=line_rgb, line_w=line_w_pt, rx=br,
                                     shape=shape_type)
            if text:
                v_anchor = "middle" if h < Inches(1.5) else "top"
                write_text_in_shape(shp, [(text, fs, color_rgb, weight, align)], anchor=v_anchor)
        else:
            # 純粋なテキスト要素 (背景・罫線なし)
            if text:
                # 高さが行数に合わせて十分かチェック
                num_lines = text.count("\n") + 1
                lh = parse_px(elem.get("lineHeight", "")) or fs * 1.5
                needed_h = lh * num_lines
                box_h = max(h, px_to_emu(needed_h, "y"))
                # 小さいフォント（フッター等）は折り返さない＆幅を少し広げる
                small_font = fs < 13
                box_w = w
                if small_font:
                    box_w = int(w * 1.15)
                tb = add_textbox(slide, x, y, box_w, box_h, "",
                                  font_size_px=fs, color_rgb=color_rgb, weight=weight, align=align)
                # 改行対応で書き直す
                tf = tb.text_frame
                tf.word_wrap = not small_font
                # クリア
                p_xml = tf._txBody
                for p_el in list(p_xml.findall(qn('a:p'))):
                    p_xml.remove(p_el)
                for line in text.split("\n"):
                    p = tf.add_paragraph()
                    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                                    "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
                    run = p.add_run()
                    run.text = line.strip()
                    run.font.size = Pt(font_size_pt(fs))
                    run.font.name = "Hiragino Sans"
                    if color_rgb:
                        run.font.color.rgb = RGBColor(*color_rgb)
                    if weight in ("700", "bold") or (str(weight).isdigit() and int(weight) >= 600):
                        run.font.bold = True
    
    # === 2. SVG要素を変換 ===
    # rect, text, line, polygon に分類
    rects = [e for e in svg_elems if e["tag"] == "rect"]
    texts = [e for e in svg_elems if e["tag"] == "text"]
    lines = [e for e in svg_elems if e["tag"] == "line"]
    polys = [e for e in svg_elems if e["tag"] == "polygon"]
    circs = [e for e in svg_elems if e["tag"] == "circle"]
    
    # rectをPPTX矩形に変換
    rect_to_shape = {}  # id(rect) → shape
    for rect in rects:
        x = px_to_emu(rect["x"], "x")
        y = px_to_emu(rect["y"], "y")
        w = px_to_emu(rect["w"], "x")
        h = px_to_emu(rect["h"], "y")
        if w <= 0 or h <= 0:
            continue
        
        fill_rgb, fill_alpha = parse_color_with_alpha(rect.get("fill", ""))
        stroke_rgb = parse_color(rect.get("stroke", "")) or None
        stroke_w = rect.get("strokeWidth", 0)
        rx = rect.get("rx", 0)
        dash = bool(rect.get("strokeDasharray", ""))
        
        if fill_rgb is None and stroke_rgb is None:
            continue
        
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 1 else MSO_SHAPE.RECTANGLE
        shp = add_filled_shape(slide, x, y, w, h,
                                fill_rgb=fill_rgb, line_rgb=stroke_rgb,
                                line_w=stroke_w, rx=rx, shape=shape_type, dash=dash)
        rect_to_shape[id(rect)] = (shp, rect)
    
    # circle (P12施策3の真円バッジなど)
    for circ in circs:
        x = px_to_emu(circ["x"], "x")
        y = px_to_emu(circ["y"], "y")
        w = px_to_emu(circ["w"], "x")
        h = px_to_emu(circ["h"], "y")
        fill_rgb = parse_color(circ.get("fill", ""))
        stroke_rgb = parse_color(circ.get("stroke", ""))
        if w > 0 and h > 0:
            shp = add_filled_shape(slide, x, y, w, h,
                                    fill_rgb=fill_rgb, line_rgb=stroke_rgb,
                                    line_w=1 if stroke_rgb else 0, shape=MSO_SHAPE.OVAL)
            rect_to_shape[id(circ)] = (shp, {
                "x": circ["x"], "y": circ["y"], "w": circ["w"], "h": circ["h"]
            })
    
    # text を rect に埋め込む or 独立配置
    # text の中心点で rect 内判定（text-anchor補正不要）
    for text in texts:
        tcx = text.get("cx", text["x"])
        tcy = text.get("cy", text["y"])
        anchor = text.get("textAnchor", "start")
        
        candidates = []
        for rid, (shp, rect_data) in rect_to_shape.items():
            if rect_data["x"] <= tcx <= rect_data["x"] + rect_data["w"] and \
               rect_data["y"] <= tcy <= rect_data["y"] + rect_data["h"]:
                candidates.append((rect_data["w"] * rect_data["h"], shp, rect_data))
        
        color_rgb = parse_color(text.get("fill", "")) or (0, 0, 0)
        fs = text.get("fontSize", 11)
        weight = text.get("fontWeight", "400")
        align = "center" if anchor == "middle" else ("right" if anchor == "end" else "left")
        
        if candidates:
            # 最も小さい (= 最も内側の) rect を選択
            candidates.sort(key=lambda c: c[0])
            _, target_shp, target_rect = candidates[0]
            # 既存のtext_frameに paragraph として追加 (累積)
            tf = target_shp.text_frame
            # 既に他のテキストがあれば改行追加
            existing_text = "".join(p.text for p in tf.paragraphs)
            if existing_text.strip():
                p = tf.add_paragraph()
            else:
                # 最初のpを使う (空のもの)
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                # 既存の空テキストrunをクリア
                for r in list(p.runs):
                    r._r.getparent().remove(r._r)
            
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                            "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
            run = p.add_run()
            run.text = text["text"]
            run.font.size = Pt(font_size_pt(fs))
            run.font.name = "Hiragino Sans"
            if color_rgb:
                run.font.color.rgb = RGBColor(*color_rgb)
            if weight in ("700", "bold") or (str(weight).isdigit() and int(weight) >= 600):
                run.font.bold = True
            
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Emu(36000)
            tf.margin_right = Emu(36000)
            tf.margin_top = Emu(18000)
            tf.margin_bottom = Emu(18000)
        else:
            # 独立テキストボックスとして配置 (bbox位置を使用)
            x_emu = px_to_emu(text["x"], "x")
            y_emu = px_to_emu(text["y"], "y")
            w_emu = px_to_emu(max(text.get("w", 50), 30), "x")
            h_emu = px_to_emu(max(text.get("h", fs * 1.3), fs * 1.3), "y")
            tb = add_textbox(slide, x_emu, y_emu, w_emu, h_emu, text["text"],
                              font_size_px=fs, color_rgb=color_rgb, weight=weight, align=align)
            tb.text_frame.word_wrap = False
    
    # line を直線シェイプに変換
    for line in lines:
        x1 = px_to_emu(line["x1"], "x")
        y1 = px_to_emu(line["y1"], "y")
        x2 = px_to_emu(line["x2"], "x")
        y2 = px_to_emu(line["y2"], "y")
        stroke_rgb = parse_color(line.get("stroke", "")) or (136, 136, 136)
        sw = line.get("strokeWidth", 1)
        dash = bool(line.get("strokeDasharray", ""))
        add_line(slide, x1, y1, x2, y2, color_rgb=stroke_rgb, width_pt=sw, dash=dash)
    
    # polygon (主に矢印先端) — bbox から RIGHT_TRIANGLE を生成
    for poly in polys:
        if "x" not in poly:
            continue
        x = px_to_emu(poly["x"], "x")
        y = px_to_emu(poly["y"], "y")
        w = px_to_emu(poly["w"], "x")
        h = px_to_emu(poly["h"], "y")
        if w <= 0 or h <= 0:
            continue
        fill_rgb = parse_color(poly.get("fill", "")) or (224, 78, 73)
        # 三角形（矢印先端）として配置
        shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(*fill_rgb)
        shp.line.fill.background()
    
    return slide


def main():
    prs = Presentation()
    prs.slide_width = DST_W
    prs.slide_height = DST_H
    
    json_files = sorted(EXTRACTED_DIR.glob("slide_*.json"))
    print(f"処理: {len(json_files)}枚")
    
    for jf in json_files:
        with open(jf, "r", encoding="utf-8") as f:
            slide_data = json.load(f)
        idx = slide_data["index"]
        print(f"  Slide {idx+1}: HTML {len(slide_data['htmlElements'])}, SVG {len(slide_data['svgElements'])}")
        build_slide(prs, slide_data)
    
    prs.save(OUTPUT_PPTX)
    size_kb = Path(OUTPUT_PPTX).stat().st_size / 1024
    print(f"\n✅ PPTX生成完了: {OUTPUT_PPTX}")
    print(f"   サイズ: {size_kb:.1f} KB")

if __name__ == "__main__":
    main()
